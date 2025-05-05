import logging
from typing import Optional

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import pymupdf
from langchain_core.tools import ToolException
from office365.runtime.auth.client_credential import ClientCredential
from office365.sharepoint.client_context import ClientContext
from pydantic import Field, PrivateAttr, create_model, model_validator, SecretStr
from transformers import BlipProcessor, BlipForConditionalGeneration

from .utils import read_docx_from_bytes
from ..elitea_base import BaseToolApiWrapper

NoInput = create_model(
    "NoInput"
)

ReadList = create_model(
    "ReadList",
    list_title=(str, Field(description="Name of a Sharepoint list to be read.")),
    limit=(Optional[int], Field(description="Limit (maximum number) of list items to be returned", default=1000))
)

GetFiles = create_model(
    "GetFiles",
    folder_name=(Optional[str], Field(description="Folder name to get list of the files.", default=None)),
    limit_files=(Optional[int], Field(description="Limit (maximum number) of files to be returned. Can be called with synonyms, such as First, Top, etc., or can be reflected just by a number for example 'Top 10 files'. Use default value if not specified in a query WITH NO EXTRA CONFIRMATION FROM A USER", default=100)),
)

ReadDocument = create_model(
    "ReadDocument",
    path=(str, Field(description="Contains the server-relative path of a document for reading.")),
    is_capture_image=(Optional[bool], Field(description="Determines is pictures in the document should be recognized.", default=False)),
    page_number=(Optional[int], Field(description="Specifies which page to read. If it is None, then full document will be read.", default=None))
)


class SharepointApiWrapper(BaseToolApiWrapper):
    site_url: str
    client_id: str = None
    client_secret: SecretStr = None
    token: SecretStr = None
    _client: Optional[ClientContext] = PrivateAttr()  # Private attribute for the office365 client

    @model_validator(mode='before')
    @classmethod
    def validate_toolkit(cls, values):
        try:
            from office365.runtime.auth.authentication_context import AuthenticationContext
            from office365.sharepoint.client_context import ClientContext
        except ImportError:
            raise ImportError(
                "`office365` package not found, please run "
               "`pip install office365-rest-python-client`"
            )

        site_url = values['site_url']
        client_id = values.get('client_id')
        client_secret = values.get('client_secret')
        token = values.get('token')

        try:
            if client_id and client_secret:
                credentials = ClientCredential(client_id, client_secret)
                cls._client = ClientContext(site_url).with_credentials(credentials)
                logging.info("Authenticated with secret id")
            elif token:
                cls._client = ClientContext(site_url).with_access_token(lambda: type('Token', (), {
                    'tokenType': 'Bearer',
                    'accessToken': token
                })())
                logging.info("Authenticated with token")
            else:
                raise ToolException("You have to define token or client id&secret.")
            logging.info("Successfully authenticated to SharePoint.")
        except Exception as e:
                logging.error(f"Failed to authenticate with SharePoint: {str(e)}")
        return values


    def read_list(self, list_title, limit: int = 1000):
        """ Reads a specified List in sharepoint site. Number of list items is limited by limit (default is 1000). """
        try:
            target_list = self._client.web.lists.get_by_title(list_title)
            self._client.load(target_list)
            self._client.execute_query()
            items = target_list.items.get().top(limit).execute_query()
            logging.info("{0} items from sharepoint loaded successfully.".format(len(items)))
            result = []
            for item in items:
                result.append(item.properties)
            return result
        except Exception as e:
            logging.error(f"Failed to load items from sharepoint: {e}")
            return ToolException("Can not list items. Please, double check List name and read permissions.")


    def get_files_list(self, folder_name: str = None, limit_files: int = 100):
        """ If folder name is specified, lists all files in this folder under Shared Documents path. If folder name is empty, lists all files under root catalog (Shared Documents). Number of files is limited by limit_files (default is 100)."""
        try:
            result = []

            target_folder_url = f"Shared Documents/{folder_name}" if folder_name else "Shared Documents"
            files = (self._client.web.get_folder_by_server_relative_path(target_folder_url)
                     .get_files(True)
                     .execute_query())

            for file in files:
                if len(result) >= limit_files:
                    break
                temp_props = {
                    'Name': file.properties['Name'],
                    'Path': file.properties['ServerRelativeUrl'],
                    'Created': file.properties['TimeCreated'],
                    'Modified': file.properties['TimeLastModified'],
                    'Link': file.properties['LinkingUrl']
                }
                result.append(temp_props)
            return result if result else ToolException("Can not get files or folder is empty. Please, double check folder name and read permissions.")
        except Exception as e:
            logging.error(f"Failed to load files from sharepoint: {e}")
            return ToolException("Can not get files. Please, double check folder name and read permissions.")

    def read_file(self, path, is_capture_image: bool = False, page_number: int = None):
        """ Reads file located at the specified server-relative path. """
        try:
            file = self._client.web.get_file_by_server_relative_path(path)
            self._client.load(file).execute_query()

            file_content = file.read()
            self._client.execute_query()
        except Exception as e:
            logging.error(f"Failed to load file from SharePoint: {e}. Path: {path}. Please, double check file name and path.")
            return ToolException("File not found. Please, check file name and path.")

        if file.name.endswith('.txt'):
            try:
                file_content_str = file_content.decode('utf-8')
            except Exception as e:
                logging.error(f"Error decoding file content: {e}")
        elif file.name.endswith('.docx'):
            file_content_str = read_docx_from_bytes(file_content)
        elif file.name.endswith('.xlsx') or file.name.endswith('.xls'):
            try:
                excel_file = io.BytesIO(file_content)
                df = pd.read_excel(excel_file)
                df.fillna('', inplace=True)
                file_content_str = df.to_string()
            except Exception as e:
                return ToolException(f"Error reading Excel file: {e}")
        elif file.name.endswith('.pdf'):
            with pymupdf.open(stream=file_content, filetype="pdf") as report:
                text_content = ''
                if page_number is not None:
                    page = report.load_page(page_number - 1)
                    text_content += self.read_pdf_page(report, page, page_number, is_capture_image)
                else:
                    for index, page in enumerate(report, start=1):
                        text_content += self.read_pdf_page(report, page, index, is_capture_image)
                file_content_str = text_content
        elif file.name.endswith('.pptx'):
            prs = Presentation(io.BytesIO(file_content))
            text_content = ''
            if page_number is not None:
                text_content += self.read_pptx_slide(prs.slides[page_number - 1], page_number, is_capture_image)
            else:
                for index, slide in enumerate(prs.slides, start=1):
                    text_content += self.read_pptx_slide(slide, index, is_capture_image)
            file_content_str = text_content
        else:
            return ToolException("Not supported type of files entered. Supported types are TXT, DOCX, PDF, PPTX, XLSX and XLS only.")
        return file_content_str

    def read_pdf_page(self, report, page, index, is_capture_images):
        text_content = f'Page: {index}\n'
        text_content += page.get_text()
        if is_capture_images:
            images = page.get_images(full=True)
            for i, img in enumerate(images):
                xref = img[0]
                base_image = report.extract_image(xref)
                img_bytes = base_image["image"]
                text_content += self.describe_image(Image.open(io.BytesIO(img_bytes)).convert("RGB"))
        return text_content

    def read_pptx_slide(self, slide, index, is_capture_image):
        text_content = f'Slide: {index}\n'
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content += shape.text + "\n"
            elif is_capture_image and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    caption = self.describe_image(Image.open(io.BytesIO(shape.image.blob)).convert("RGB"))
                except:
                    caption = "\n[Picture: unknown]\n"
                text_content += caption
        return text_content

    def describe_image(self, image):
        processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
        model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
        inputs = processor(image, return_tensors="pt")
        out = model.generate(**inputs)
        return "\n[Picture: " + processor.decode(out[0], skip_special_tokens=True) + "]\n"

    def get_available_tools(self):
        return [
            {
                "name": "read_list",
                "description": self.read_list.__doc__,
                "args_schema": ReadList,
                "ref": self.read_list
            },
            {
                "name": "get_files_list",
                "description": self.get_files_list.__doc__,
                "args_schema": GetFiles,
                "ref": self.get_files_list
            },
            {
                "name": "read_document",
                "description": self.read_file.__doc__,
                "args_schema": ReadDocument,
                "ref": self.read_file
            }
        ]