from docx import Document
from io import BytesIO
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import pymupdf
from langchain_core.tools import ToolException
from transformers import BlipProcessor, BlipForConditionalGeneration

def parse_file_content(file_name, file_content, is_capture_image: bool = False, page_number: int = None):
    if file_name.endswith('.txt'):
        return parse_txt(file_content)
    elif file_name.endswith('.docx'):
        return read_docx_from_bytes(file_content)
    elif file_name.endswith('.xlsx') or file_name.endswith('.xls'):
        return parse_excel(file_content)
    elif file_name.endswith('.pdf'):
        return parse_pdf(file_content, page_number, is_capture_image)
    elif file_name.endswith('.pptx'):
        return parse_pptx(file_content, page_number, is_capture_image)
    else:
        return ToolException(
            "Not supported type of files entered. Supported types are TXT, DOCX, PDF, PPTX, XLSX and XLS only.")

def parse_txt(file_content):
    try:
        return file_content.decode('utf-8')
    except Exception as e:
        return ToolException(f"Error decoding file content: {e}")

def parse_excel(file_content):
    try:
        excel_file = io.BytesIO(file_content)
        df = pd.read_excel(excel_file)
        df.fillna('', inplace=True)
        return df.to_string()
    except Exception as e:
        return ToolException(f"Error reading Excel file: {e}")

def parse_pdf(file_content, page_number, is_capture_image):
    with pymupdf.open(stream=file_content, filetype="pdf") as report:
        text_content = ''
        if page_number is not None:
            page = report.load_page(page_number - 1)
            text_content += read_pdf_page(report, page, page_number, is_capture_image)
        else:
            for index, page in enumerate(report, start=1):
                text_content += read_pdf_page(report, page, index, is_capture_image)
        return text_content

def parse_pptx(file_content, page_number, is_capture_image):
    prs = Presentation(io.BytesIO(file_content))
    text_content = ''
    if page_number is not None:
        text_content += read_pptx_slide(prs.slides[page_number - 1], page_number, is_capture_image)
    else:
        for index, slide in enumerate(prs.slides, start=1):
            text_content += read_pptx_slide(slide, index, is_capture_image)
    return text_content

def read_pdf_page(report, page, index, is_capture_images):
    text_content = f'Page: {index}\n'
    text_content += page.get_text()
    if is_capture_images:
        images = page.get_images(full=True)
        for i, img in enumerate(images):
            xref = img[0]
            base_image = report.extract_image(xref)
            img_bytes = base_image["image"]
            text_content += describe_image(Image.open(io.BytesIO(img_bytes)).convert("RGB"))
    return text_content

def read_docx_from_bytes(file_content):
    """Read and return content from a .docx file using a byte stream."""
    try:
        doc = Document(BytesIO(file_content))
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return '\n'.join(text)
    except Exception as e:
        print(f"Error reading .docx from bytes: {e}")
        return ""

def read_pptx_slide(slide, index, is_capture_image):
    text_content = f'Slide: {index}\n'
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text_content += shape.text + "\n"
        elif is_capture_image and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                caption = describe_image(Image.open(io.BytesIO(shape.image.blob)).convert("RGB"))
            except:
                caption = "\n[Picture: unknown]\n"
            text_content += caption
    return text_content

def describe_image(image):
    processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
    model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
    inputs = processor(image, return_tensors="pt")
    out = model.generate(**inputs)
    return "\n[Picture: " + processor.decode(out[0], skip_special_tokens=True) + "]\n"
